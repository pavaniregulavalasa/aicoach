"""
Document Generator - Creates PDF and PPT documents from training content
"""
import os
import logging
from pathlib import Path
from datetime import datetime
from typing import Optional
import re

logger = logging.getLogger(__name__)

# Try to import PDF libraries
try:
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False
    logger.warning("reportlab not available - PDF generation disabled")

# Try to import PPT libraries
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available - PPT generation disabled")


def markdown_to_paragraphs(text: str, styles):
    """Convert markdown text to ReportLab Paragraph objects"""
    paragraphs = []
    lines = text.split('\n')
    current_paragraph = []
    current_level = 0
    
    for line in lines:
        line = line.strip()
        if not line:
            if current_paragraph:
                paragraphs.append(Paragraph(' '.join(current_paragraph), styles['Normal']))
                paragraphs.append(Spacer(1, 0.1*inch))
                current_paragraph = []
            continue
        
        # Check for headers
        if line.startswith('##'):
            if current_paragraph:
                paragraphs.append(Paragraph(' '.join(current_paragraph), styles['Normal']))
                current_paragraph = []
            header_text = line.lstrip('#').strip()
            paragraphs.append(Paragraph(header_text, styles['Heading2']))
            paragraphs.append(Spacer(1, 0.2*inch))
        elif line.startswith('###'):
            if current_paragraph:
                paragraphs.append(Paragraph(' '.join(current_paragraph), styles['Normal']))
                current_paragraph = []
            header_text = line.lstrip('#').strip()
            paragraphs.append(Paragraph(header_text, styles['Heading3']))
            paragraphs.append(Spacer(1, 0.15*inch))
        elif line.startswith('#'):
            if current_paragraph:
                paragraphs.append(Paragraph(' '.join(current_paragraph), styles['Normal']))
                current_paragraph = []
            header_text = line.lstrip('#').strip()
            paragraphs.append(Paragraph(header_text, styles['Heading1']))
            paragraphs.append(Spacer(1, 0.3*inch))
        # Check for bullet points
        elif line.startswith('- ') or line.startswith('* '):
            if current_paragraph:
                paragraphs.append(Paragraph(' '.join(current_paragraph), styles['Normal']))
                current_paragraph = []
            bullet_text = line[2:].strip()
            paragraphs.append(Paragraph(f"• {bullet_text}", styles['Normal']))
        # Check for numbered lists
        elif re.match(r'^\d+\.\s', line):
            if current_paragraph:
                paragraphs.append(Paragraph(' '.join(current_paragraph), styles['Normal']))
                current_paragraph = []
            list_text = re.sub(r'^\d+\.\s', '', line)
            paragraphs.append(Paragraph(f"• {list_text}", styles['Normal']))
        else:
            current_paragraph.append(line)
    
    if current_paragraph:
        paragraphs.append(Paragraph(' '.join(current_paragraph), styles['Normal']))
    
    return paragraphs


def generate_pdf(training_content: str, title: str, level: str, knowledge_base: str, output_path: str) -> str:
    """
    Generate a PDF document from training content
    
    Args:
        training_content: Markdown-formatted training content
        title: Document title
        level: Training level
        knowledge_base: Knowledge base name
        output_path: Path to save the PDF
        
    Returns:
        Path to generated PDF file
    """
    if not REPORTLAB_AVAILABLE:
        raise ImportError("reportlab library is not installed. Install it with: pip install reportlab")
    
    logger.info(f"Generating PDF: {output_path}")
    
    # Create output directory if it doesn't exist
    output_dir = os.path.dirname(output_path)
    if output_dir:
        os.makedirs(output_dir, exist_ok=True)
    
    # Create PDF document
    doc = SimpleDocTemplate(output_path, pagesize=letter,
                          rightMargin=72, leftMargin=72,
                          topMargin=72, bottomMargin=18)
    
    # Container for the 'Flowable' objects
    story = []
    
    # Define styles
    styles = getSampleStyleSheet()
    
    # Custom styles
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        textColor=colors.HexColor('#1f4788'),
        spaceAfter=30,
        alignment=TA_CENTER,
        fontName='Helvetica-Bold'
    )
    
    subtitle_style = ParagraphStyle(
        'CustomSubtitle',
        parent=styles['Normal'],
        fontSize=12,
        textColor=colors.HexColor('#666666'),
        spaceAfter=20,
        alignment=TA_CENTER
    )
    
    header_style = ParagraphStyle(
        'CustomHeader',
        parent=styles['Heading2'],
        fontSize=16,
        textColor=colors.HexColor('#1f4788'),
        spaceAfter=12,
        spaceBefore=12,
        fontName='Helvetica-Bold'
    )
    
    subheader_style = ParagraphStyle(
        'CustomSubHeader',
        parent=styles['Heading3'],
        fontSize=14,
        textColor=colors.HexColor('#2c5aa0'),
        spaceAfter=8,
        spaceBefore=8,
        fontName='Helvetica-Bold'
    )
    
    # Add title page
    story.append(Paragraph(title, title_style))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph(f"Knowledge Base: {knowledge_base.upper()}", subtitle_style))
    story.append(Paragraph(f"Level: {level.upper()}", subtitle_style))
    story.append(Paragraph(f"Generated: {datetime.now().strftime('%B %d, %Y')}", subtitle_style))
    story.append(PageBreak())
    
    # Process markdown content
    paragraphs = markdown_to_paragraphs(training_content, {
        'Heading1': title_style,
        'Heading2': header_style,
        'Heading3': subheader_style,
        'Normal': styles['Normal']
    })
    
    story.extend(paragraphs)
    
    # Build PDF
    doc.build(story)
    
    logger.info(f"PDF generated successfully: {output_path}")
    return output_path


def generate_ppt(training_content: str, title: str, level: str, knowledge_base: str, output_path: str) -> str:
    """
    Generate a PowerPoint presentation from training content
    
    Args:
        training_content: Markdown-formatted training content
        title: Presentation title
        level: Training level
        knowledge_base: Knowledge base name
        output_path: Path to save the PPT
        
    Returns:
        Path to generated PPT file
    """
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx library is not installed. Install it with: pip install python-pptx")
    
    logger.info(f"Generating PPT: {output_path}")
    
    # Create output directory if it doesn't exist
    output_dir = os.path.dirname(output_path)
    if output_dir:
        os.makedirs(output_dir, exist_ok=True)
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = f"{knowledge_base.upper()} - {level.upper()}\n{datetime.now().strftime('%B %d, %Y')}"
    
    # Set title formatting
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 72, 136)
    
    # Process content into slides
    lines = training_content.split('\n')
    current_slide_content = []
    current_slide_title = None
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
        
        # Check for main headers (##) - new slide
        if line.startswith('##'):
            # Save previous slide if exists
            if current_slide_title or current_slide_content:
                _add_content_slide(prs, current_slide_title, current_slide_content)
            
            # Start new slide
            current_slide_title = line.lstrip('#').strip()
            current_slide_content = []
        # Check for subheaders (###) - section within slide
        elif line.startswith('###'):
            if current_slide_content:
                current_slide_content.append(('subheader', line.lstrip('#').strip()))
        # Check for bullet points
        elif line.startswith('- ') or line.startswith('* '):
            bullet_text = line[2:].strip()
            current_slide_content.append(('bullet', bullet_text))
        # Check for numbered lists
        elif re.match(r'^\d+\.\s', line):
            list_text = re.sub(r'^\d+\.\s', '', line)
            current_slide_content.append(('bullet', list_text))
        # Regular text
        else:
            if line and not line.startswith('#'):
                current_slide_content.append(('text', line))
    
    # Add last slide
    if current_slide_title or current_slide_content:
        _add_content_slide(prs, current_slide_title, current_slide_content)
    
    # Save presentation
    prs.save(output_path)
    
    logger.info(f"PPT generated successfully: {output_path}")
    return output_path


def _add_content_slide(prs, title, content):
    """Helper function to add a content slide to presentation"""
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Add title if exists
    if title:
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(0.8)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 72, 136)
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    # Add content
    if content:
        left = Inches(0.5)
        top = Inches(1.5) if title else Inches(0.5)
        width = Inches(9)
        height = Inches(5.5) if title else Inches(6.5)
        content_box = slide.shapes.add_textbox(left, top, width, height)
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for item_type, item_text in content:
            p = content_frame.add_paragraph()
            p.text = item_text
            
            if item_type == 'subheader':
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = RGBColor(44, 90, 160)
                p.space_after = Pt(6)
            elif item_type == 'bullet':
                p.level = 0
                p.font.size = Pt(14)
                p.space_after = Pt(4)
                # Add bullet point
                p.text = f"• {item_text}"
            else:  # text
                p.font.size = Pt(12)
                p.space_after = Pt(4)


def generate_document(training_content: str, title: str, level: str, knowledge_base: str, 
                     format_type: str = "pdf", output_dir: str = "generated_docs") -> Optional[str]:
    """
    Generate a document (PDF or PPT) from training content
    
    Args:
        training_content: Markdown-formatted training content
        title: Document title
        level: Training level
        knowledge_base: Knowledge base name
        format_type: "pdf" or "ppt"
        output_dir: Directory to save the document
        
    Returns:
        Path to generated document or None if generation failed
    """
    try:
        # Ensure output_dir is relative to project root (AICOACH folder)
        # Get the project root directory (where this file is located)
        project_root = Path(__file__).parent.parent
        output_dir_path = project_root / output_dir
        
        # Create safe filename
        safe_title = re.sub(r'[^\w\s-]', '', title).strip()
        safe_title = re.sub(r'[-\s]+', '-', safe_title)
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"{safe_title}_{level}_{timestamp}.{format_type.lower()}"
        output_path = str(output_dir_path / filename)
        
        if format_type.lower() == "pdf":
            if not REPORTLAB_AVAILABLE:
                logger.error("reportlab not available - cannot generate PDF")
                return None
            return generate_pdf(training_content, title, level, knowledge_base, output_path)
        elif format_type.lower() in ["ppt", "pptx"]:
            if not PPTX_AVAILABLE:
                logger.error("python-pptx not available - cannot generate PPT")
                return None
            # Always use .pptx extension
            output_path = output_path.replace('.ppt', '.pptx')
            return generate_ppt(training_content, title, level, knowledge_base, output_path)
        else:
            logger.error(f"Unsupported format: {format_type}")
            return None
    except Exception as e:
        logger.error(f"Error generating {format_type} document: {str(e)}")
        logger.exception("Full traceback:")
        return None

